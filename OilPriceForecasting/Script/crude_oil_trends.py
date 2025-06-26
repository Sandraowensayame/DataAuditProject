from pptx import Presentation
from pptx.util import Inches

# Create Presentation
prs = Presentation()
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# Add Title and Subtitle
title = slide.shapes.title
title.text = "Crude Oil Prices: Trends & Forecast (1987–2034)"
subtitle = slide.placeholders[1]
subtitle.text = "Historical Insights, Market Dynamics, and Projected Outlook"

# Add Main Text
slide = prs.slides.add_slide(prs.slide_layouts[5])
text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
tf = text_box.text_frame
tf.text = "Key Insights:"
points = [
    "• Prices stable until early 2000s, surged mid-decade, and became volatile post-2020.",
    "• Gradual recovery trend observed post-2020 pandemic.",
    "• Forecast suggests a moderate rise until 2034, aiding in strategic decisions."
]
for point in points:
    p = tf.add_paragraph()
    p.text = point
    p.level = 1

# Save Presentation
pptx_file = 'crude_oil_trends_forecast.pptx'
prs.save(pptx_file)
